from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pandas import DataFrame as DF

def set_reverse_categories(axis):
    """
    workaround function that replicates the "Categories in Reverse Order" UI option in PPT
    """
    ele = axis._element.xpath(r'c:scaling/c:orientation')[0]
    ele.set("val", "maxMin")

p = Presentation()
# Create some data to be used in the chart
series_names = ["A","B","C","D"]
cat_names = ["cat 1"]
data = {
        cat_names[0]: [.10, .20, .30, .40]
        }
df = DF(data, series_names, cat_names)
cd = CategoryChartData()
cd.categories = df.index
for name in df.columns:
    data = df[name]
    cd.add_series(name, data, '0%')

layout = p.slide_layouts[6] # MODIFY AS NEEDED, 6 is the index of my "Blank" slide template.

# Create two charts, one reversed and one not reversed on the Category Axis
for reverse in (True, False):
    slide = p.slides.add_slide( layout )
    shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, 0, 0, 9143301, 6158000, cd) 
    cht = shape.chart
    plot = cht.plots[0]
    plot.has_data_labels = False
    if reverse:
        set_reverse_categories(cht.category_axis)

p.save('/Users/milosdjelic/Desktop/QBR/example.pptx')
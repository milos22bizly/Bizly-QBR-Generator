
import pandas as pd
import numpy as np
import mysql.connector
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Cm,Pt
from pptx import chart
from pptx.enum.chart import XL_LEGEND_POSITION,XL_LABEL_POSITION
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE
from pandas import DataFrame as DF
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION

""" ****************************************************************** CONECTION TO DB- MYSQL *************************************************************************"""
mydb = mysql.connector.connect(
    host="bizly-db-prod-reader-do-user-12143907-0.b.db.ondigitalocean.com",
    database='defaultdb',
    user="qbr-reader",
    password="AVNS__uHZWnkHP0xksBmbFcL",
    port=25060  
    )
mydb.connect(db='bizly_prod')

mydb2 = mysql.connector.connect(
host="bizly-dev-do-user-12143907-0.b.db.ondigitalocean.com",
database='bizly_dev_20220801',
user="doadmin",
password="AVNS_JeUb_97kvCdsOz8uun-",
port=25060  
)


sql_query1="""
select distinct
	"teams"."id" AS "team_id",
	"teams"."name" AS "company_name",
    "users"."id" AS "id",
    year("users"."created_at") AS "year",
    month("users"."created_at") AS "month"
from ("users" join "teams" on(("teams"."id" = "users"."current_team_id")) ) 
where "users"."deleted_at" is Null 
union select 'team_id',
'ALL TEAMS' AS "company_name",
"users"."id" AS "id",
year("users"."created_at") AS "year",
month("users"."created_at") AS "month"
from "users" 
order by "year" desc,"month" desc
"""
mycursor=mydb.cursor()
mycursor.execute(sql_query1)
myresult = mycursor.fetchall()


registred_users=pd.DataFrame(myresult,columns=['team_id','company_name','id','year','month'])
registred_users['company_name']=registred_users['company_name'].astype('string')
registred_users['team_id']=registred_users['team_id'].astype('string')

sql_query1="""
    select 
        "teams"."id" AS "team_id",
        year("meetings"."created_at") AS "year",
        month("meetings"."created_at") AS "month",
        count(0) AS "meetings_created",
        count("booking_inquiries"."id") AS "meetings_with_inquiries" 
    from (("meetings" join "teams" on(("teams"."id" = "meetings"."team_id"))) 
        left join "booking_inquiries" on(("booking_inquiries"."meeting_id" = "meetings"."id"))) 
        where ("teams"."name" <> 'Bizly') 
        group by "teams"."id",month("meetings"."created_at"),
        year("meetings"."created_at") union select 0 AS "team_id",
        year("meetings"."created_at") AS "year",month("meetings"."created_at") AS "month",
        count(0) AS "meetings_created",count("booking_inquiries"."id") AS "meetings_with_inquiries" 
    from (("meetings" left join "booking_inquiries" on(("booking_inquiries"."meeting_id" = "meetings"."id"))) 
        join "teams" on(("teams"."id" = "meetings"."team_id"))) 
        where ("teams"."name" <> 'Bizly') 
        group by month("meetings"."created_at"),
        year("meetings"."created_at") 
    union select "teams"."id" AS "team_id",
    0 AS "year",0 AS "month",
    count(0) AS "meetings_created",
    count("booking_inquiries"."id") AS "meetings_with_inquiries" 

    from (("meetings" join "teams" on(("teams"."id" = "meetings"."team_id"))) 
    left join "booking_inquiries" on(("booking_inquiries"."meeting_id" = "meetings"."id"))) 
    where ("teams"."name" <> 'Bizly') 
    group by "teams"."id" union select 0 AS "team_id",
    0 AS "year",
    0 AS "month",count(0) AS "meetings_created",count("booking_inquiries"."id") AS "meetings_with_inquiries" 
    from (("meetings" join "teams" on(("teams"."id" = "meetings"."team_id"))) 
    left join "booking_inquiries" on(("booking_inquiries"."meeting_id" = "meetings"."id"))) 
    where ("teams"."name" <> 'Bizly') order by "year" desc,"month" desc;
"""
mycursor=mydb.cursor()
mycursor.execute(sql_query1)
myresult = mycursor.fetchall()


team_meeting=pd.DataFrame(myresult,columns=['team_id','year','month','meetings_created','meetings_with_inquiries'])
team_meeting['team_id']=team_meeting['team_id'].astype('string')

data_all=pd.merge(registred_users,team_meeting, how='left' ,left_on=['team_id','month','year'], right_on = ['team_id','month','year'])
data_all=data_all.fillna(0)

sql_query2="""
select
	"proposals"."team_id",
    year("proposals"."accepted_at") AS "year",
    month("proposals"."accepted_at") AS "month",
    u.id as user_id
from "proposals" left join users u on proposals.accepted_by = u.id
where proposals.deleted_at is null 
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query2)
myresult = mycursor.fetchall()

accepted_user=pd.DataFrame(myresult,columns=['team_id','year','month','user_id'])
accepted_user['team_id']=accepted_user['team_id'].astype('string')

sql_query2="""
select convert("orders"."team_id", char) as team_id,
       user_id,
    year("orders"."contracted_at") AS "year",
    month("orders"."contracted_at") AS "month"
from orders 
where orders.deleted_at is null 
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query2)
myresult = mycursor.fetchall()

bookings_confirmed=pd.DataFrame(myresult,columns=['team_id','user_id','year','month'])
bookings_confirmed['team_id']=bookings_confirmed['team_id'].astype('string')

sql_query_meetings1="""
select  
	"meetings"."team_id", 
    year("meetings"."created_at") AS "year",
    month("meetings"."created_at") AS "month",
    count("meetings"."created_at") as Created
from "meetings"
where meetings.deleted_at is null and meetings.starts_at is not null and meetings.ends_at is not null
group by "meetings"."team_id",month("meetings"."created_at"),year("meetings"."created_at")
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_meetings1)
myresult = mycursor.fetchall()

meetings=pd.DataFrame(myresult,columns=['team_id','year','month','Created'])
meetings['team_id']=meetings['team_id'].astype('string')

sql_query_meetings2="""
select
    "orders"."team_id",
    year("orders".final_spend_tracked_at) AS "year",
    month("orders"."final_spend_tracked_at") AS "month",
    count("orders"."final_spend_tracked_at") as Completed
from "orders"
where orders.deleted_at is null and orders.final_spend_tracked_at is not null and orders.final_spend_tracked_at is not null
group by "orders"."team_id",month("orders"."final_spend_tracked_at"),year("orders"."final_spend_tracked_at")
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_meetings2)
myresult = mycursor.fetchall()

meetings2=pd.DataFrame(myresult,columns=['team_id','year','month','Completed'])
meetings2['team_id']=meetings2['team_id'].astype('string')

meetings=pd.merge(meetings,meetings2, how='left' ,left_on=['team_id','month','year'], right_on = ['team_id','month','year'])

sql_query_meetings3="""
select
    "meetings"."team_id",
    year("meetings"."cancelled_at") AS "year",
    month("meetings"."cancelled_at") AS "month",
    count("meetings"."cancelled_at") as Cancelled
from "meetings"
where meetings.deleted_at is null and meetings.starts_at is not null and meetings.ends_at is not null and meetings.demo_mode=FALSE
group by "meetings"."team_id",month("meetings"."cancelled_at"),year("meetings"."cancelled_at")
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_meetings3)
myresult = mycursor.fetchall()

meetings3=pd.DataFrame(myresult,columns=['team_id','year','month','Cancelled'])
meetings3['team_id']=meetings3['team_id'].astype('string')

meetings=pd.merge(meetings,meetings3, how='left' ,left_on=['team_id','month','year'], right_on = ['team_id','month','year'])


sql_query_proposals="""
select  
	"proposals"."team_id", 
    year("proposals"."accepted_at") AS "year",
    month("proposals"."accepted_at") AS "month",
    count("proposals"."accepted_at") as Accepted
from "proposals"
where proposals.deleted_at is null
group by "proposals"."team_id",month("proposals"."accepted_at"),year("proposals"."accepted_at")  
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_proposals)
myresult = mycursor.fetchall()

proposals=pd.DataFrame(myresult,columns=['team_id','year','month','Accepted'])
proposals['team_id']=proposals['team_id'].astype('string')


meetings=pd.merge(meetings,proposals, how='left' ,left_on=['team_id','month','year'], right_on = ['team_id','month','year'])

sql_query_bookings="""
select
    "booking_inquiries"."team_id",
    year("booking_inquiries"."submitted_at") AS "year",
    month("booking_inquiries"."submitted_at") AS "month",
    count("booking_inquiries"."submitted_at") as Submitted
from "booking_inquiries" join meetings m on booking_inquiries.meeting_id = m.id
where booking_inquiries.deleted_at is  null and booking_inquiries.demo_mode=FALSE and booking_inquiries.submitted_at is not null and m.deleted_at is null and m.cancelled_at is null
and m.starts_at is not null and m.ends_at is not null
group by "booking_inquiries"."team_id",month("booking_inquiries"."submitted_at"),year("booking_inquiries"."submitted_at")

"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_bookings)
myresult = mycursor.fetchall()

bookings=pd.DataFrame(myresult,columns=['team_id','year','month','Submitted'])
bookings['team_id']=bookings['team_id'].astype('string')

meetings=pd.merge(meetings,bookings, how='left' ,left_on=['team_id','month','year'], right_on = ['team_id','month','year'])

sql_query_orders="""

select  
	convert("orders"."team_id", char) as team_id,
    year("orders"."contracted_at") AS "year",
    month("orders"."contracted_at") AS "month",
    count("orders"."contracted_at") as Booked
from "orders"
where orders.deleted_at is null
group by "orders"."team_id",month("orders"."contracted_at"),year("orders"."contracted_at")

"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_orders)
myresult = mycursor.fetchall()

orders=pd.DataFrame(myresult,columns=['team_id','year','month','Booked'])
orders['team_id']=orders['team_id'].astype('string')

sql_query_booked="""

select  
	convert("orders"."team_id", char) as team_id,
    year("orders"."contracted_at") AS "year",
    month("orders"."contracted_at") AS "month",
    "orders"."proposal_id" as proposal_id
from "orders"
where orders.deleted_at is null
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_booked)
myresult = mycursor.fetchall()

booked_meetings=pd.DataFrame(myresult,columns=['team_id','year','month','proposal_id'])
booked_meetings['team_id']=booked_meetings['team_id'].astype('string')
booked_meetings['month']=booked_meetings['month'].fillna(0)
booked_meetings['year']=booked_meetings['year'].fillna(0)
booked_meetings['month']=booked_meetings['month'].astype('int64')
booked_meetings['year']=booked_meetings['year'].astype('int64')

meetings=pd.merge(meetings,orders, how='left' ,left_on=['team_id','month','year'], right_on = ['team_id','month','year'])
meetings=meetings.fillna(0)

sql_avg_lead_time="""
select 
meetings.team_id,
month(meetings.created_at) as month,
year(meetings.created_at) as year,
avg(TIMESTAMPDIFF(week,meetings.created_at,meetings.starts_at_utc))  as avg_lead_time
from meetings
where meetings.deleted_at is NULL
and meetings.created_at is not NULL
and meetings.starts_at_utc is not NULL and meetings.starts_at is not null and meetings.ends_at is not null
group by meetings.team_id,month(meetings.created_at),year(meetings.created_at)

"""
mycursor=mydb.cursor()
mycursor.execute(sql_avg_lead_time)
myresult = mycursor.fetchall()

avg_lead_time=pd.DataFrame(myresult,columns=['team_id','month','year','avg_lead_time_calculated'])
avg_lead_time['team_id']=avg_lead_time['team_id'].astype('string')

sql_avg_resp_time="""
select * from 
	(select booking_inquiry_venues.created_by, 
			year(booking_inquiry_venues.updated_at) as year,
			month(booking_inquiry_venues.updated_at) as month,
			avg(booking_inquiry_venues.response_time) as response_time
            from booking_inquiry_venues
            where booking_inquiry_venues.response_time is not null
            group by booking_inquiry_venues.created_by, year(booking_inquiry_venues.updated_at), month(booking_inquiry_venues.updated_at) 
            ) as t1
		left join
	(select users.id,
    	    convert(users.current_team_id, char) as team_id
        from users)
        as t2
	on t1.created_by=t2.id
"""
mycursor=mydb.cursor()
mycursor.execute(sql_avg_resp_time)
myresult = mycursor.fetchall()

avg_resp_time=pd.DataFrame(myresult,columns=['created_by','year','month','response_time','id','team_id'])
avg_resp_time['team_id']=avg_resp_time['team_id'].astype('string')


sql_query_pvr="""
select team_id,month, 
		year, 
        Sum(proposal_count) AS proposal_count,
       Sum(sent_count + rejected_count) AS response_count,
       ( Sum(sent_count + rejected_count) / Sum(proposal_count) ) AS response_rate
 from (
SELECT 
	p.meeting_id  AS meeting_id,
    p.team_id as team_id,
   Min(Month(p.created_at)) AS month,
   Min(Year(p.created_at)) AS year,
   Count(p.id) AS proposal_count
    FROM   proposals AS p
        GROUP  BY p.meeting_id,p.team_id) as t1
left join 
(
SELECT meeting_id,Count(*) as sent_count
                FROM   proposals AS p2
                WHERE  sent_at IS NOT NULL 
                group by meeting_id
) as t2
on t1.meeting_id=t2.meeting_id
left join
(
	select meeting_id, count(reject) as rejected_count from (
SELECT id, meeting_id FROM  proposals AS p
) as t4
left join 
(SELECT proposal_id, Count(*) as reject
                FROM   booking_inquiry_venues AS bivs
                WHERE  bivs.rejected_at IS NOT NULL 
                group by proposal_id) as t5
on t4.id=t5.proposal_id
group by meeting_id
) as t3
on t1.meeting_id=t3.meeting_id
GROUP  BY team_id,month, year
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_pvr)
myresult = mycursor.fetchall()

pvr=pd.DataFrame(myresult,columns=['team_id','month','year','proposal_count','response_count','response_rate'])

pvr['team_id']=pvr['team_id'].astype('string')

pvr['month']=pvr['month'].fillna(0)
pvr['year']=pvr['year'].fillna(0)


sql_query_venues_types="""
select * from
	(select
    "properties"."id",
    "orders"."proposal_id",
    "orders"."team_id",
	"orders"."venue_id",
    year("orders"."confirmed_at") AS "year",
    month("orders"."confirmed_at") AS "month"
from ("properties" left join "orders" on("properties"."id" = "orders"."venue_id"))
where "orders"."contracted_at" is not null) as t1 left join
(select
	"properties"."property_type_id",
    "properties"."id",
	"property_types"."name"
from ("properties" left join "property_types" on("properties"."property_type_id" = "property_types"."id"))) as t2
on t1.venue_id=t2.id
"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_venues_types)
myresult = mycursor.fetchall()

venue_types=pd.DataFrame(myresult,columns=['id','proposal_id','team_id','venue_id','year','month','property_type_id','id2','venue_name'])
venue_types['team_id']=venue_types['team_id'].astype('string')
venue_types['month']=venue_types['month'].fillna(0)
venue_types['year']=venue_types['year'].fillna(0)

sql_contracts="""
select * from(
select 
	"booking_inquiry_venues"."id" AS "booking_inquiry_venue_id",
    "booking_inquiry_venues"."proposal_id" AS "booking_inquiry_proposal_id",
    "booking_inquiry_venues"."sent_at" AS "date_sent",
    "booking_inquiry_venues"."property_id" AS "venue_id",
    month("booking_inquiry_venues"."created_at") as month,
    year("booking_inquiry_venues"."created_at") as year,
    "booking_inquiries"."meeting_id",
    booking_inquiry_venues.preferred_venue
    from ("booking_inquiry_venues" join "booking_inquiries" on(("booking_inquiries"."id" = "booking_inquiry_venues"."inquiry_id")))
    where ("booking_inquiry_venues"."sent_at" is not null  )
    ) as t1 left join
    (select "meetings"."id", "meetings"."type", convert("meetings"."team_id",char) as team_id from "meetings") as t2
    on t1.meeting_id=t2.id 
    left join
    (select "proposals"."id" as id2, "proposals"."cancelled_at", "proposals"."sent_at", "proposals"."property_id" from "proposals") as t3
    on t1.booking_inquiry_proposal_id=t3.id2
    left join
    (select
	properties.property_type_id,
    properties.id as properties_id,
	property_types.name
    from (properties left join property_types on(properties.property_type_id = property_types.id))) as t4
    on t3.property_id=t4.properties_id
"""
mycursor=mydb.cursor()
mycursor.execute(sql_contracts)
myresult = mycursor.fetchall()

contracts=pd.DataFrame(myresult,columns=['id','booking_inquiry_proposal_id','date','venue_id','month','year','meeting_id','is_prefered','id2','type','team_id','id3','cancelled_at','sent_at','p_id','p_type','p_id2','pro_type'])
contracts['team_id']=contracts['team_id'].astype('string')
contracts['month']=contracts['month'].fillna(0)
contracts['year']=contracts['year'].fillna(0)
contracts['pnp']=np.where(contracts['is_prefered']==1,'Preferred','Non-Preferred')

contracts=contracts[contracts['pro_type']=='Hotel']

sql_query_meeting_type="""
select distinct * from(
select  
	"proposals"."team_id", 
    "proposals"."id",
    "proposals"."property_id", 
    "proposals"."created_at",
    year("proposals"."created_at") AS "year",
    month("proposals"."created_at") AS "month"
from "proposals") as t1 left join (select "properties"."property_type_id","properties"."id" from "properties") as t2
on t1.property_id=t2.id
left join (select "proposal_meeting_rooms"."proposal_id" as meeting_room from "proposal_meeting_rooms") as t3
on t1.id=t3.meeting_room
left join (select "proposal_guest_rooms"."proposal_id" as guest_room from "proposal_guest_rooms") as t4
on t1.id=t4.guest_room

"""

mycursor=mydb.cursor()
mycursor.execute(sql_query_meeting_type)
myresult = mycursor.fetchall()

meeting_type=pd.DataFrame(myresult,columns=['team_id','id','property_id','created_at','year','month','property_type_id','id2','meeting_room','guest_room'])
meeting_type['team_id']=meeting_type['team_id'].astype('string')

conditions = [
    (meeting_type['property_type_id'] == 2),
    (meeting_type['meeting_room'] > 0) & (meeting_type['guest_room'] > 0),
    (meeting_type['meeting_room'].isnull()) & (meeting_type['guest_room'] > 0),
    (meeting_type['meeting_room']>0) & (meeting_type['guest_room'].isnull()) 
    ]

values = ['Restaurant', 'Meeting Space & Guest Rooms', 'Guest Room Only', 'Meeting Room']

meeting_type['mt'] = np.select(conditions, values)

meeting_type['mt']=np.where(meeting_type['mt']=="0",'Other',meeting_type['mt'])

sql_location="""
select * from (
	select 
convert("meetings"."team_id",char) AS "team_id",
"meetings"."id",
"meetings"."status_id",
"meetings"."city_id" AS "city_id",
"meetings"."deleted_at" as deleted_at
 from meetings ) as t1 left join
(select "cities"."id","cities"."state", "cities"."name" from "cities") as t2
on t1.city_id=t2.id
left join (
select "orders"."event_id",
year("orders"."contracted_at") AS "year",
month("orders"."contracted_at") AS "month"
from orders
) as t3 
on t1.id=t3.event_id
where t1.status_id in ('1','2','3') and t1.deleted_at is null
"""

mycursor=mydb.cursor()
mycursor.execute(sql_location)
myresult = mycursor.fetchall()

cities=pd.DataFrame(myresult,columns=['team_id','id','status_id','city_id','deleted_at','id2','state','name','event_id','year','month'])
cities['team_id']=cities['team_id'].astype('string')


""" ****************************************************************** CONECTION TO DB- MYSQL *************************************************************************"""

""" ****************************************************************** TRANSFORMATION OF RAW DATA  *************************************************************************"""

list_clients=['10982', '13449', '10983', '802', '419'] 
#list_clients=['419']
""" 
months=[1,2,3]
q_end=3
list_month=['January','February','March']
q_text="Q1"
first_month_period=['2022-01-01','2022-01-31']
second_month_period=['2022-02-01','2022-02-28']
third_month_period=['2022-03-01','2022-03-31']


months=[4,5,6]
q_end=6
list_month=['April','May','June']
q_text="Q2"
first_month_period=['2022-04-01','2022-04-30']
second_month_period=['2022-05-01','2022-05-31']
third_month_period=['2022-06-01','2022-06-30']

months=[7,8,9]
q_end=9
list_month=['July','August','September']
q_text="Q3"
first_month_period=['2022-07-01','2022-07-31']
second_month_period=['2022-08-01','2022-08-31']
third_month_period=['2022-09-01','2022-09-30']
"""
year=2022
months_=(10,11,12)
months=[10,11,12]
q_end=12
list_month=['October','November','December']
q_text="Q4"
first_month_period=['2022-10-01','2022-10-31']
second_month_period=['2022-11-01','2022-11-30']
third_month_period=['2022-12-01','2022-12-31']


""" ****************************************************************** TRANSFORMATION OF RAW DATA  *************************************************************************"""


""" ****************************************************************** CREATION OF PPTX SLIDE 1  *************************************************************************"""

for client in list_clients:

    print("*********")
    print(client)
    print("*********")

    total_registrations=registred_users[registred_users['team_id']==client]
    tot_reg=total_registrations['id'].nunique()

    sql_query_contr_final="""
    select distinct * from(
    select
        `proposals`.`team_id`,
        `proposals`.`id`,
        `proposals`.`property_id`,
        `proposals`.`accepted_at`,
        `proposals`.`meeting_id`
    from `proposals`) as t1 left join (select `properties`.`property_type_id`,`properties`.`id` from `properties`) as t2
    on t1.property_id=t2.id
    left join (select
            `proposal_meeting_rooms`.`proposal_id` as meeting_room
            from `proposal_meeting_rooms`
            ) as t3
    on t1.id=t3.meeting_room
    left join (select
        `proposal_guest_rooms`.`proposal_id` as guest_room,
        `proposal_guest_rooms`.`rate` as guest_room_rate,
        `proposal_guest_rooms`.`quantity` as guest_room_quantity,
        proposal_guest_rooms.rate * proposal_guest_rooms.quantity as `guest_room_total`
        from `proposal_guest_rooms`) as t4
    on t1.id=t4.guest_room
    left join(select
    `orders`.`proposal_id` as proposal_id,
    `orders`.`confirmed_at` as confirmed_at,
    `orders`.`fb_spend` as fb_spend,
    `orders`.`av_spend` as av_spend,
    `orders`.`room_spend` as room_spend,
    `orders`.`prod_spend` as prod_spend,
    `orders`.`transaction_amount` as transaction_amount,
    `orders`.`total_est` as total_est,
    year(`orders`.`confirmed_at`) AS `year`,
    month(`orders`.`confirmed_at`) AS `month`,
    `orders`.`meeting_room_est` as meeting_room_rate,
    `orders`.`fb_est` as fb_minimum,
    `orders`.`guestroom_est` as guestroom_rate,

    `orders`.`attendee_est` as guests
    from `orders`) as t5
    on t1.id=t5.proposal_id
    left join (select
    `meetings`.`id` as meeting_ind,
    `meetings`.`cancelled_at`
    from `meetings`) as t6
    on t1.meeting_id=t6.meeting_ind
    where t1.team_id='{c}'
    and t1.accepted_at is not null and t6.cancelled_at is null
    and t5.confirmed_at is not null
    """.format(c=client)

    mycursor=mydb.cursor()
    mycursor.execute(sql_query_contr_final)
    myresult = mycursor.fetchall()

    contr_final=pd.DataFrame(myresult,columns=['team_id',	'id',	'property_id',	'accepted_at',	'meeting_id',	'property_type_id',	'id22',	'meeting_room',	
    'guest_room',	'guest_room_rate',	'guest_room_quantity',	'guest_room_total',	'proposal_id',	'confirmed_at',	'fb_spend',	'av_spend',	'room_spend',	'prod_spend',	
    'transaction_amount','total_est','year',	'month','meeting_room_rate', 'fb_minimum','guestroom_rate','guests',	'meeting_ind',	'cancelled_at'])


    contr_final['team_id']=contr_final['team_id'].astype('string')
    contr_final['guest_room_rate']=contr_final['guest_room_rate'].astype('float')
    contr_final['guest_room_quantity']=contr_final['guest_room_quantity'].astype('float')
    contr_final['meeting_room_rate']=contr_final['meeting_room_rate'].astype('float')
    contr_final['fb_minimum']=contr_final['fb_minimum'].astype('float')
    contr_final['total_est']=contr_final['total_est'].astype('float')
    contr_final['transaction_amount']=contr_final['transaction_amount'].astype('float')


    
    contr_final['guest_room_rate']=contr_final['guest_room_rate'].fillna(0)
    contr_final['guest_room_quantity']=contr_final['guest_room_quantity'].fillna(0)
    contr_final['meeting_room_rate']=contr_final['meeting_room_rate'].fillna(0)
    contr_final['meeting_room_rate']=contr_final['meeting_room_rate'].fillna(0)
    contr_final['guestroom_rate']=contr_final['guestroom_rate'].fillna(0)

    contr_final['fb_minimum']=contr_final['fb_minimum'].fillna(0)
    contr_final['fb_spend']=contr_final['fb_spend'].fillna(0)
    contr_final['av_spend']=contr_final['av_spend'].fillna(0)
    contr_final['room_spend']=contr_final['room_spend'].fillna(0)
    contr_final['prod_spend']=contr_final['prod_spend'].fillna(0)
    contr_final['total_est']=contr_final['total_est'].fillna(0)
    contr_final['transaction_amount']=contr_final['transaction_amount'].fillna(0)

    #contr_final['contracted_value']=contr_final['guest_room_rate']*contr_final['guest_room_quantity'] + contr_final['meeting_room_rate'] + contr_final['fb_minimum']
    #contr_final['final_value']=contr_final['fb_spend']+contr_final['av_spend'] + contr_final['room_spend'] + contr_final['prod_spend']

    contr_final['contracted_value']=contr_final['total_est']
    contr_final['final_value']=contr_final['transaction_amount']


    sql_query_con_values="""
        select distinct * from(
    select
        `proposals`.`team_id`,
        `proposals`.`id`,
        `proposals`.`property_id`,
        `proposals`.`accepted_at`
    from `proposals`) as t1 left join (select `properties`.`property_type_id`,`properties`.`id` from `properties`) as t2
    on t1.property_id=t2.id

    left join(select
    `orders`.`proposal_id` as proposal_id,
    `orders`.`confirmed_at` as confirmed_at,
    `orders`.`transaction_amount` as transaction_amount,
    `orders`.`total_est` as total_est,
    year(`orders`.`confirmed_at`) AS `year`,
    month(`orders`.`confirmed_at`) AS `month`
    from `orders`) as t5
    on t1.id=t5.proposal_id
    where t1.team_id={c} 
    and t1.accepted_at is not null
    and t5.confirmed_at is not null
    """.format(c=client)

    mycursor=mydb.cursor()
    mycursor.execute(sql_query_con_values)
    myresult = mycursor.fetchall()

    con_value=pd.DataFrame(myresult,columns=['team_id',	'id',	'property_id',	'accepted_at',	'property_type_id',	'id',
    'proposal_id',	'confirmed_at','transaction_amount','total_est','year',	'month'])

    con_value['team_id']=con_value['team_id'].astype('string')
    con_value['total_est']=con_value['total_est'].astype('float')
    con_value['transaction_amount']=con_value['transaction_amount'].astype('float')
    con_value['total_est']=con_value['total_est'].fillna(0)
    con_value['transaction_amount']=con_value['transaction_amount'].fillna(0)



    sql_query_hotel="""
    SELECT
    property_brands.name as 'Brand',
    property_brands.id as 'id',

    property_types.name as 'Type',
    YEAR(orders.contracted_at) as 'Year',
    MONTH(orders.contracted_at) as 'Month',
    orders.team_id as 'team_id',
    count(orders.id) as 'number_of_meetings',
    orders.proposal_id,
    SUM(orders.transaction_amount) as 'Total Booking Amount',
    SUM(orders.total_est) as 'Total Bookings (Estimated)',
    sum(orders.guestroom_spend) as 'GuestRoomTotal'
        
    FROM orders
    left join properties on properties.id = orders.venue_id
    left join property_chains on properties.chain_id = property_chains.id
    left join property_brands on property_chains.brand_id = property_brands.id
    left join property_types on(properties.property_type_id = property_types.id)
    WHERE orders.contracted_at is not NULL and orders.deleted_at is null and YEAR(orders.contracted_at)='2022' and orders.team_id='{c}'

    GROUP BY 
    property_brands.name,
    property_types.name,
    MONTH(orders.contracted_at), 
    YEAR(orders.contracted_at),
    orders.team_id, orders.proposal_id,
    property_brands.id


    ORDER BY YEAR(orders.contracted_at),
    MONTH(orders.contracted_at)
    """.format(c=client)

    mycursor=mydb.cursor()
    mycursor.execute(sql_query_hotel)
    myresult = mycursor.fetchall()

    hotel_brands=pd.DataFrame(myresult,columns=['brand','id','p_type','year','month','team_id','number_of_meetings','proposal_id','total_booking_amount','total_bookings_estimated','guestroomtotal'])
    hotel_brands['team_id']=hotel_brands['team_id'].astype('string')

    hotel_brands['len']=hotel_brands['brand'].str.len()

    hotel_brands['brand']=np.where(hotel_brands['len']>1,hotel_brands['brand'],'Independent')

    #hotel_brands=hotel_brands[hotel_brands['id']!=32]


    active_user1="""
    select id from `users` where `current_team_id` ='{c}'  and `email` not like '%bizly%'
    and exists (select * from `meetings` where `users`.`id` = `meetings`.`created_by` and `created_at` <= '{end}'
    and `created_at` >= '{start}' and `name` not like '%test%' 
    and `meetings`.`deleted_at` is null and `meetings`.`starts_at` is not null 
    and `meetings`.`ends_at` is not null
    and not exists (select * from `event_templates` 
    where `meetings`.`id` = `event_templates`.`event_id`)) and `users`.`deleted_at` is null
    """.format(c=client,end=first_month_period[1], start=first_month_period[0])
    mycursor=mydb.cursor()
    mycursor.execute(active_user1)
    myresult = mycursor.fetchall()

    active_user1=pd.DataFrame.from_records(myresult, columns =['ids'])


    active_user2="""
    select id from `users` where `current_team_id` = '{c}'  and `email` not like '%bizly%'
    and exists (select * from `meetings` where `users`.`id` = `meetings`.`created_by` and `created_at` <= '{end}'
    and `created_at` >= '{start}' and `name` not like '%test%' 
    and `meetings`.`deleted_at` is null
    and `meetings`.`starts_at` is not null 
    and `meetings`.`ends_at` is not null
    and not exists (select * from `event_templates` 
    where `meetings`.`id` = `event_templates`.`event_id`)) and `users`.`deleted_at` is null
    """.format(c=client,end=second_month_period[1], start=second_month_period[0])

    mycursor=mydb.cursor()
    mycursor.execute(active_user2)
    myresult = mycursor.fetchall()

    active_user2=pd.DataFrame.from_records(myresult, columns =['ids'])

    active_user3="""
    select id from `users` where `current_team_id` = '{c}'  and `email` not like '%bizly%'
    and exists (select * from `meetings` where `users`.`id` = `meetings`.`created_by` and `created_at` <= '{end}'
    and `created_at` >= '{start}' and `name` not like '%test%' 
    and `meetings`.`deleted_at` is null and `meetings`.`starts_at` is not null 
    and `meetings`.`ends_at` is not null
    and not exists (select * from `event_templates` 
    where `meetings`.`id` = `event_templates`.`event_id`)) and `users`.`deleted_at` is null
    """.format(c=client,end=third_month_period[1], start=third_month_period[0])

    mycursor=mydb.cursor()
    mycursor.execute(active_user3)
    myresult = mycursor.fetchall()
    
    active_user3=pd.DataFrame.from_records(myresult, columns =['ids'])


    all_users=pd.concat([active_user1, active_user2,active_user3])
    all_users=all_users.drop_duplicates(subset=['ids'])
    all_users=all_users.shape[0]

    active_user_m1=active_user1.shape[0]
    active_user_m2=active_user2.shape[0]
    active_user_m3=active_user3.shape[0]

    playbook="""
    SELECT q1.playbook,
       q1.playbook_type,
       
       CASE
         WHEN playbook_type = 'Global' THEN 'Bizly Concierge'
         WHEN playbook_type = 'Team' THEN teams.name
         WHEN playbook_type = 'User' THEN Concat(u.first_name, ' ', u.last_name)
       END AS Creator,
       q1.use_count
FROM   (SELECT event_templates.name                              AS playbook,
               IF(event_templates.team_id IS NULL
                  AND event_templates.user_id IS NULL, 'Global', IF(
               event_templates.team_id IS NOT NULL
               AND
               event_templates.user_id IS NULL, 'Team', 'User')) AS
               playbook_type,
               event_templates.user_id,
               event_templates.team_id,
               Count(meetings.id)                                AS use_count
        FROM   meetings
               join event_templates
                 ON event_templates.id = meetings.template_id
        WHERE  meetings.deleted_at IS NULL
               AND meetings.team_id = '{c}'
               AND Month(meetings.created_at) IN {m}
               AND Year(meetings.created_at) = {y} and meetings.id in (select id from(
                                                                    select  meetings.id,
                                                                        meetings.team_id, 
                                                                        year(meetings.created_at) AS year,
                                                                        month(meetings.created_at) AS month,
                                                                        count(meetings.created_at) as Created
                                                                    from meetings
                                                                    where meetings.deleted_at is null and meetings.starts_at is not null and meetings.ends_at is not null
                                                                    and meetings.team_id='{c}' and year(meetings.created_at)={y} and month(meetings.created_at) in {m}
                                                                    group by meetings.id, meetings.team_id,month(meetings.created_at),year(meetings.created_at)) as t1)
        GROUP  BY playbook,
                  playbook_type,
                  event_templates.user_id,
                  event_templates.team_id
        ORDER  BY use_count DESC) AS q1
       left join users u
              ON q1.user_id = u.id
       left join teams
              ON q1.team_id = teams.id
WHERE  ( playbook_type = 'User'
         AND u.first_name IS NOT NULL )
        OR playbook_type <> 'User' and playbook not in ('Blank')
    """.format(c=client,y=year,m=months_)

    mycursor=mydb.cursor()
    mycursor.execute(playbook)
    myresult = mycursor.fetchall()
    
    team_creator_client=pd.DataFrame.from_records(myresult, columns =['playbook','playbook_type','Creator','count'])

    team_creator_client=team_creator_client.drop(columns='playbook_type')

    
    data=registred_users[registred_users['team_id']==client]
    bookings_confirmed_client=bookings_confirmed[bookings_confirmed['team_id']==client]
    accepted_user_client=accepted_user[accepted_user['team_id']==client]
    pvr_client=pvr[pvr['team_id']==client]
    venue_types_client=venue_types[venue_types['team_id']==client]
    contracts_client=contracts[contracts['team_id']==client]
    meeting_type_client=meeting_type[meeting_type['team_id']==client]
    contr_final_client=contr_final[contr_final['team_id']==client]
    con_value_client=con_value[con_value['team_id']==client]

    cities_client=cities[cities['team_id']==client]
    hotel_brands_client=hotel_brands[hotel_brands['team_id']==client]
    booked_meetings_client=booked_meetings[booked_meetings['team_id']==client]

    client_name=str(data['company_name'].unique()[0])
    print(client_name)
    
    data=data[data['year']==2022]
    bookings_confirmed_client=bookings_confirmed_client[bookings_confirmed_client['year']==2022]
    accepted_user_client=accepted_user_client[accepted_user_client['year']==2022]

    pvr_client=pvr_client[pvr_client['year']==2022]
    venue_types_client=venue_types_client[venue_types_client['year']==2022]
    contracts_client=contracts_client[contracts_client['year']==2022]
    #meeting_type_client=meeting_type_client[meeting_type_client['year']==2022]
    contr_final_client=contr_final_client[contr_final_client['year']==2022]
    con_value_client=con_value_client[con_value_client['year']==2022]

    cities_client=cities_client[cities_client['year']==2022]
    hotel_brands_client=hotel_brands_client[hotel_brands_client['year']==2022]
    booked_meetings_client=booked_meetings_client[booked_meetings_client['year']==2022]

    data=data[data['month'].isin(months)]
    pvr_client=pvr_client[pvr_client['month'].isin(months)]
    venue_types_client=venue_types_client[venue_types_client['month'].isin(months)]
    booked_meetings_client=booked_meetings_client[booked_meetings_client['month'].isin(months)]

    #contracts_client=contracts_client[contracts_client['month'].isin(months)]

    contracts_client=contracts_client[contracts_client['booking_inquiry_proposal_id'].isin(booked_meetings_client['proposal_id'].values.tolist())]
    venue_types_client=venue_types_client[venue_types_client['proposal_id'].isin(booked_meetings_client['proposal_id'].values.tolist())]

    con_value_client=con_value_client[con_value_client['proposal_id'].isin(hotel_brands_client['proposal_id'].values.tolist())]
    contr_final_client=contr_final_client[contr_final_client['proposal_id'].isin(hotel_brands_client['proposal_id'].values.tolist())]

    venue_types_client=venue_types_client[['venue_name','team_id']]
    venue_types_client=venue_types_client.groupby(by=["venue_name"]).count().reset_index()

    

    contracts_client=contracts_client[['pnp','type']]
    contracts_client_internal=contracts_client[contracts_client['type']=='Internal']
    contracts_client_external=contracts_client[contracts_client['type']=='External']
    

    contracts_client=contracts_client.groupby(by=["pnp"]).count().reset_index()
    contracts_client_internal=contracts_client_internal.groupby(by=["pnp"]).count().reset_index()
    contracts_client_external=contracts_client_external.groupby(by=["pnp"]).count().reset_index()

    #meeting_type_client=meeting_type_client[meeting_type_client['month'].isin(months)]
    #meeting_type_client=meeting_type_client[meeting_type_client['mt'].isin(values)]

    meeting_type_client=meeting_type_client[meeting_type_client['id'].isin(booked_meetings_client['proposal_id'].values.tolist())]

    meeting_type_client=meeting_type_client.groupby(by=["mt"]).count().reset_index()

    #con_value_client['total']=con_value_client['guest_room_rate']*con_value_client['guest_room_quantity'] + con_value_client['meeting_room_rate'] + con_value_client['fb_minimum']
    

    contr_final_client=contr_final_client[contr_final_client['month'].isin(months)]
    con_value_client=con_value_client[con_value_client['month'].isin(months)]
    con_value_client=con_value_client[['total_est','transaction_amount','month']]
    
    con_value_client=con_value_client.groupby(by=["month"]).sum().reset_index()

    con_value_client_org=contr_final_client

    cities_client=cities_client[cities_client['month'].isin(months)]

    cities_client=cities_client[['state','name','id']]
    cities_client_state=cities_client.groupby(by=["state","name"]).count().reset_index()

    hotel_brands_client=hotel_brands_client[hotel_brands_client['month'].isin(months)]


    book_brand_hotel=hotel_brands_client[hotel_brands_client['p_type']=='Hotel']

    book_brand_hotel=book_brand_hotel[['brand','number_of_meetings']]
    book_brand_hotel=book_brand_hotel.groupby(by=["brand"]).sum().reset_index()
    book_brand_hotel['id']=book_brand_hotel['number_of_meetings']

    book_brand_non_hotel=hotel_brands_client[hotel_brands_client['p_type']!='Hotel']

    book_brand_non_hotel=book_brand_non_hotel[['p_type','number_of_meetings']]
    book_brand_non_hotel=book_brand_non_hotel.groupby(by=["p_type"]).sum().reset_index()
    book_brand_non_hotel['id']=book_brand_non_hotel['number_of_meetings']
    
    
    #hotel_brands_client['total']=np.where(hotel_brands_client['total_bookings_estimated']>0,hotel_brands_client['total_bookings_estimated'],hotel_brands_client['total_booking_amount'])
    hotel_brands_client.to_excel('aaa.xlsx')
    hotel_brands_client['total']=hotel_brands_client['total_bookings_estimated']
    #hotel_brands_client['total']=np.where(hotel_brands_client['total_booking_amount']>0,hotel_brands_client['total_booking_amount'],hotel_brands_client['total_bookings_estimated'])

    hotel_brands_client['total']=hotel_brands_client['total'].fillna(0)
    #hotel_brands_client['total']=hotel_brands_client['total'].astype(int)

    brand_value_hotel=hotel_brands_client[hotel_brands_client['p_type']=='Hotel']
    brand_value_hotel=brand_value_hotel[['brand','total']]
    brand_value_hotel=brand_value_hotel[brand_value_hotel['total']>0]


    brand_value_non_hotel=hotel_brands_client[hotel_brands_client['p_type']!='Hotel']
    brand_value_non_hotel=brand_value_non_hotel[['p_type','total']]
    brand_value_non_hotel=brand_value_non_hotel[brand_value_non_hotel['total']>0]

    if brand_value_hotel.shape[0]>0:
        brand_value_hotel=brand_value_hotel.groupby(by=["brand"]).sum().reset_index()

    if brand_value_non_hotel.shape[0]>0:
        brand_value_non_hotel=brand_value_non_hotel.groupby(by=["p_type"]).sum().reset_index()
        
        #brand_value['total']=round(brand_value['total'],0).astype(float)

    
    bookings_confirmed_client=bookings_confirmed_client[bookings_confirmed_client['month'].isin(months)]
    accepted_user_client=accepted_user_client[accepted_user_client['month'].isin(months)]

    data=data.sort_values(by='month',ascending=True)

    prs = Presentation('/Users/milosdjelic/Desktop/QBR/Template.pptx')

    slide = prs.slides[0]
    
    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{} 2022 Reporting".format(q_text)


    font = run.font
    font.name = 'Courier New'
    font.size = Pt(18)
    font.color.rgb = RGBColor(255,255,255)

    image=client+".png"
    x, y, cx, cy = Cm(1.25), Cm(12.75), Cm(14.6), Cm(5.2)
    slide.shapes.add_picture(image, x, y, cx, cy)

    slide = prs.slides[1]   

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "User Behavior {}".format(q_text)

    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)

    chart_data = CategoryChartData()
    chart_data.categories = list_month
   

    chart_data.add_series('Registered Users', [data[data['month']==months[0]]['id'].count(),data[data['month']==months[1]]['id'].count(),data[data['month']==months[2]]['id'].count()])
    chart_data.add_series('Active Users', [active_user_m1,active_user_m2,active_user_m3])

    x, y, cx, cy = Cm(0.5), Cm(4), Cm(16), Cm(10)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(18)
    data_labels.bold=True
    data_labels.font.color.rgb = RGBColor(0,0,0)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END

    left = Cm(6) 
    width = height = Cm(3)
    top=Cm(2.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = "Quarterly Users"    
    p.font.bold = True		 

    def set_reverse_categories(axis):
        ele = axis._element.xpath(r'c:scaling/c:orientation')[0]
        ele.set("val", "maxMin")

    # Create some data to be used in the chart
    series_names = ["Registered Users","Active Users","Accepted Users","Booked Users"]
    registred=data['id'].count()
    active=all_users
   
    try:
        accepted=len(accepted_user_client['user_id'].unique())
        accepted=int(accepted)
    except:
        accepted=0
    try:
        booked=len(bookings_confirmed_client['user_id'].unique())
        booked=int(booked)
    except:
        booked=0

    cat_names = ["Quarterly Unique Users"]
    data = {
            cat_names[0]: [registred, active, accepted, booked]
            }
    df = DF(data, series_names, cat_names)
    
    cd = CategoryChartData()
    cd.categories = df.index
    for name in df.columns:
        data = df[name]
        cd.add_series(name, data)

    #layout = prs.slide_layouts[1] # MODIFY AS NEEDED, 6 is the index of my "Blank" slide template.
    slide = prs.slides[1]   

    x, y, cx, cy = Cm(17), Cm(3), Cm(16), Cm(10)
    # Create two charts, one reversed and one not reversed on the Category Axis
    #for reverse in (True, False):
    #slide = prs.slides.add_slide( layout )
    shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, cd) 
    chart = shape.chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(16)
    data_labels.bold=True
    data_labels.font.color.rgb = RGBColor(255,255,255)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END


    #if reverse:
    set_reverse_categories(chart.category_axis)
    axis = chart.category_axis
    #axis.reverse_order = False

    left = Cm(24) 
    width = height = Cm(3)
    top=Cm(15.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = str(tot_reg)   
    p.font.bold = True


    ppt_name=client_name+"_"+q_text+'.pptx'
    
    meetings_client=meetings[meetings['team_id']==client]
    meetings_client=meetings_client[meetings_client['year']==2022]
    meetungs_client_ytd=meetings_client
    meetings_client=meetings_client[meetings_client['month'].isin(months)]
    meetings_client=meetings_client.sort_values(by='month',ascending=True)

    lead_time_client=avg_lead_time[avg_lead_time['team_id']==client]
    lead_time_client=lead_time_client[lead_time_client['year']==2022]
    lead_time_client=lead_time_client[lead_time_client['month'].isin(months)]

    avg_lead_time_calculated=lead_time_client['avg_lead_time_calculated'].sum() / lead_time_client['avg_lead_time_calculated'].count() 


    resp_time_client=avg_resp_time[avg_resp_time['team_id']==client]
    resp_time_client=resp_time_client[resp_time_client['year']==2022]
    resp_time_client=resp_time_client[resp_time_client['month'].isin(months)]
    resp_time_client=resp_time_client.sort_values(by='month',ascending=True)

    resp_time=resp_time_client['response_time'].sum() / resp_time_client['response_time'].count() / 8

    pvr_value=pvr_client['response_rate'].mean()

    slide = prs.slides[2]

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Meetings Summary for {}".format(q_text)


    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)

 
    chart_data = CategoryChartData()
    chart_data.categories = list_month
    chart_data.add_series('Events Created', meetings_client['Created'])
    chart_data.add_series('Events Inquired', meetings_client['Submitted'])
    chart_data.add_series('Events Accepted', meetings_client['Accepted'])
    chart_data.add_series('Events Booked', meetings_client['Booked'])
    chart_data.add_series('Events Completed', meetings_client['Completed'])
    chart_data.add_series('Events Cancelled', meetings_client['Cancelled'])

    

    x, y, cx, cy = Cm(0.5), Cm(4), Cm(16), Cm(10)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(16)
    data_labels.bold=True
    data_labels.font.color.rgb = RGBColor(0,0,0)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END

    left = Cm(6) 
    width = height = Cm(3)
    top=Cm(2.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = "Quarterly Meeting Activity"    
    p.font.bold = True	

    left = Cm(1.8) 
    width = height = Cm(3)
    top=Cm(14)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    left = Cm(23.31) 
    width = height = Cm(2)
    top=Cm(2.65)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    p = tf.add_paragraph()		                          
    p.text = "Average Meeting Lead Time - {} Weeks".format(str(int(avg_lead_time_calculated)))    
    p.font.bold = True	

    p = tf.add_paragraph()		                          
    p.text = "Average Venue Response Time - {} Business Days".format(str(float(round(resp_time,2))))
    p.font.bold = True	

    p = tf.add_paragraph()		                          
    p.text = "Percent Venue Response - {:.2%}".format(pvr_value)      
    p.font.bold = True	

    
    
    left = Cm(23.31) 
    width = height = Cm(2)
    top=Cm(2.65)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame 
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetings_client['Created'].sum())))
    p.font.bold = True	

    left = Cm(23.51) 
    width = height = Cm(2)
    top=Cm(3.43)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetings_client['Submitted'].sum()))) 
    p.font.bold = True	

    
    left = Cm(23.74) 
    width = height = Cm(2)
    top=Cm(4.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetings_client['Accepted'].sum())))   
    p.font.bold = True	

    
    left = Cm(23.34) 
    width = height = Cm(2)
    top=Cm(4.92)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetings_client['Booked'].sum())))
    p.font.bold = True	

    
    left = Cm(24.14) 
    width = height = Cm(2)
    top=Cm(5.67)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetings_client['Completed'].sum())))
    p.font.bold = True	

    
    
    left = Cm(23.7) 
    width = height = Cm(2)
    top=Cm(6.47)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetings_client['Cancelled'].sum())))
    p.font.bold = True	
    
    left = Cm(23.96) 
    width = height = Cm(2)
    top=Cm(10.53)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetungs_client_ytd['Created'].sum())))    
    p.font.bold = True

    left = Cm(24.14) 
    width = height = Cm(2)
    top=Cm(11.36)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetungs_client_ytd['Submitted'].sum()))) 
    p.font.bold = True	
    
    
    
    left = Cm(24.38) 
    width = height = Cm(2)
    top=Cm(12.11)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetungs_client_ytd['Accepted'].sum())))   
    p.font.bold = True	

    left = Cm(23.96) 
    width = height = Cm(2)
    top=Cm(12.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetungs_client_ytd['Booked'].sum())))
    p.font.bold = True	


    

    left = Cm(24.8) 
    width = height = Cm(2)
    top=Cm(13.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetungs_client_ytd['Completed'].sum())))
    p.font.bold = True	

    

    

    left = Cm(24.17) 
    width = height = Cm(2)
    top=Cm(14.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()		                          
    p.text = "{}".format(str(int(meetungs_client_ytd['Cancelled'].sum())))
    p.font.bold = True	
    
    
    

    slide = prs.slides[3]

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Meetings Summary for {}".format(q_text)


    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)

    chart_data = ChartData()

    if venue_types_client.shape[0]>0:

        if venue_types_client.shape[0]==1:
            vt=[1]
        else:
            vt=venue_types_client['team_id']/venue_types_client['team_id'].sum()
        
        chart_data.categories =venue_types_client['venue_name'] +" - "+ venue_types_client['team_id'].astype(str)

        chart_data.add_series('Venue Type Breakdown', vt)

        x, y, cx, cy = Cm(2.13), Cm(9.34), Cm(13), Cm(8)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    chart_data = ChartData()

    if contracts_client.shape[0]>0:

        contracts_client=contracts_client.sort_values(by='pnp', ascending=False)

        chart_data = ChartData()
        chart_data.categories =contracts_client['pnp']

        chart_data.add_series('', contracts_client['type']/contracts_client['type'].sum())

        x, y, cx, cy = Cm(-0.9), Cm(2.43), Cm(8), Cm(6)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    if contracts_client_internal.shape[0]>0:

        contracts_client_internal=contracts_client_internal.sort_values(by='pnp', ascending=False)

        chart_data = ChartData()
        chart_data.categories =contracts_client['pnp']

        chart_data.add_series('',contracts_client_internal['type']/contracts_client_internal['type'].sum())

        x, y, cx, cy = Cm(4.38), Cm(1.46), Cm(8.5), Cm(7)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    if contracts_client_external.shape[0]>0:

        contracts_client_external=contracts_client_external.sort_values(by='pnp', ascending=False)

        chart_data = ChartData()
        chart_data.categories =contracts_client_external['pnp']
        chart_data.add_series('',contracts_client_external['type']/contracts_client_external['type'].sum())

        x, y, cx, cy = Cm(9.94), Cm(2.4), Cm(8), Cm(6)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    left = Cm(0.5) 
    height = Cm(1.8)
    width = Cm(10)
    top=Cm(7.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = "Total Contracts {}".format(str(contracts_client['type'].sum()))

    left = Cm(5.87) 
    height = Cm(1.8)
    width = Cm(10)
    top=Cm(7.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = "Internal Total {}".format(str(contracts_client_internal['type'].sum()))

    left = Cm(11.3) 
    height = Cm(1.8)
    width = Cm(10)
    top=Cm(7.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = "External Total {}".format(str(contracts_client_external['type'].sum()))

    chart_data = ChartData()
    
    chart_data.categories =meeting_type_client['mt'].astype('string') +' - '+ meeting_type_client['team_id'].astype('string')

    chart_data.add_series('Meeting Type Breakdown',(meeting_type_client['id']/meeting_type_client['id'].sum()))

    x, y, cx, cy = Cm(17), Cm(1.8), Cm(16), Cm(16)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart


    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    slide = prs.slides[4]   

    try:
        m1_value=con_value_client[con_value_client['month']==months[0]]
        m1_value=int(m1_value['total_est'].values[0])
    except:
        m1_value=0

    try:
        m2_value=con_value_client[con_value_client['month']==months[1]]
        m2_value=int(m2_value['total_est'].values[0])

    except:
        m2_value=0

    try:
        m3_value=con_value_client[con_value_client['month']==months[2]]
        m3_value=int(m3_value['total_est'].values[0])

    except:
        m3_value=0

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Meetings Summary for {}".format(q_text)

    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)


    chart_data = CategoryChartData()
    chart_data.categories = list_month
    chart_data.add_series('Total Contracted per Month',[m1_value,m2_value,m3_value])

    x, y, cx, cy = Cm(0.5), Cm(2), Cm(16), Cm(10)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"#,##0'

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(18)
    data_labels.bold=True
    data_labels.font.color.rgb = RGBColor(0,0,0)
    data_labels.number_format = '"$"#,##0'
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    
    left = Cm(8.7) 
    height = Cm(1.8)
    width = Cm(3)
    top=Cm(12.75)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    
    tf = txBox.text_frame
    text_1=con_value_client_org[con_value_client_org['meeting_room_rate']>0]
    text_1=round(text_1['meeting_room_rate'].mean(),3)
    text_1=float(text_1)

    p = tf.add_paragraph()		                          
    p.text = "${}".format(str((f'{text_1:,.2f}'))) 

    left = Cm(8.1) 
    height = Cm(1.8)
    width = Cm(3)
    top=Cm(13.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)
    
    tf = txBox.text_frame
    
    text_2=con_value_client_org['guestroom_rate'].sum()
    text_2=float(text_2)
    p = tf.add_paragraph()		                          
    p.text = "${}".format(str((f'{text_2:,.2f}')))

    left = Cm(4.9) 
    height = Cm(1.8)
    width = Cm(3)
    top=Cm(14.25)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()	
    guests=con_value_client_org[con_value_client_org['guests']>0]   
    guests=guests['guests']

    text_3=guests.mean()
    text_3=float(text_3)

    p.text = "{}".format(str(f'{text_3:,.0f}'))

    left = Cm(7.2) 
    height = Cm(1.8)
    width = Cm(3)
    top=Cm(15)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    con_value_client_org['total3']=con_value_client_org['guest_room_rate']*con_value_client_org['guest_room_quantity']
    con_value_client_org=con_value_client_org[con_value_client_org['total3']>0]
    con_value_client_org=con_value_client_org['total3']

    p = tf.add_paragraph()		              

    text_4=(text_1+text_2) / text_3          
    p.text = "${}".format(str(f'{text_4:,.2f}'))
    
    # Create some data to be used in the chart
    series_names = ["Contracted Value","Final Value"]
    contracted_value=con_value_client['total_est'].sum()
    fin_value=con_value_client['transaction_amount'].sum()
    contracted_value=int(contracted_value)
    
    fin_value=int(fin_value)
   
    cat_names = ["Contracted vs Final Totals*"]
    data = {
            cat_names[0]: [contracted_value, fin_value]
            }
    df = DF(data, series_names, cat_names)
    df.sort_values(by="Contracted vs Final Totals*", ascending=False)

    cd = CategoryChartData()
    cd.categories = df.index
    for name in df.columns:
        data = df[name]
        cd.add_series(name, data)

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Meetings Summary for {}".format(q_text)


    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)


    x, y, cx, cy = Cm(17), Cm(2), Cm(16), Cm(10)
  
    shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, cd) 
    chart = shape.chart

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.font.size = Pt(16)
    data_labels.bold=True
    data_labels.font.color.rgb = RGBColor(255,255,255)
    data_labels.number_format = '"$"#,##0'
    data_labels.position = XL_LABEL_POSITION.INSIDE_END

    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '"$"#,##0'

    #if reverse:
    set_reverse_categories(chart.category_axis)
    axis = chart.category_axis
    #axis.reverse_order = False




    slide = prs.slides[5]

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Location Summary for {}".format(q_text)

    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)


    if cities_client_state.shape[0]>0:
        max_rows=min(12,cities_client_state.shape[0]+1)
        x, y, cx, cy = Cm(0.5), Cm(3.31), Cm(16), Cm(10)
        shape = slide.shapes.add_table(max_rows,3, x, y, cx, cy)
        shape.has_table
        table = shape.table
        cell = table.cell(0, 0)
        cell.text = 'State'
        cell = table.cell(0, 1)
        cell.text = 'City'
        cell = table.cell(0, 2)
        cell.text = 'Number of Meetings'

        cities_client_state=cities_client_state.sort_values(by="id", ascending=False)

        for row_ in range(1,max_rows,1):
            for col_ in range(3):
                cell = table.cell(row_, col_)
                cell.text= str(cities_client_state.iloc[row_-1,col_])

    if team_creator_client.shape[0]>0:
        max_rows=min(5,team_creator_client.shape[0]+1)
        x, y, cx, cy = Cm(17.55), Cm(3.22), Cm(15.82), Cm(8.74)
        shape = slide.shapes.add_table(max_rows,3, x, y, cx, cy)
        shape.has_table
        table = shape.table
        cell = table.cell(0, 0)
        cell.text = 'Playbook'
        cell = table.cell(0, 1)
        cell.text = 'Creator'
        cell = table.cell(0, 2)
        cell.text = 'Meetings'

        team_creator_client=team_creator_client.sort_values(by="count", ascending=False)

        for row_ in range(1,max_rows,1):
            for col_ in range(3):
                cell = table.cell(row_, col_)
                cell.text= str(team_creator_client.iloc[row_-1,col_])

    slide = prs.slides[6]

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame


    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Hotel Summary for {}".format(q_text)

    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)

    if book_brand_hotel.shape[0]>0:
        book_brand_hotel['brand_text']=book_brand_hotel['brand'] + " - " + book_brand_hotel['id'].astype(str)

        book_brand_hotel=book_brand_hotel.sort_values(by='brand_text',ascending=True)

        chart_data = ChartData()
        chart_data.categories =book_brand_hotel['brand_text']

        chart_data.add_series('Booked Brand Distribution Hotels', book_brand_hotel['id']/book_brand_hotel['id'].sum())

        x, y, cx, cy = Cm(0.5), Cm(2), Cm(15), Cm(16)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart


        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'

        number_of_colors = book_brand_hotel.shape[0]
        import random

        color_list=[]
        color_list = [''.join([random.choice('0123456789ABCDEF') for j in range(6)]) for i in range(number_of_colors)]
        book_brand_list=book_brand_hotel['brand'].unique()

        try:
            Hilton_index=book_brand_list.tolist().index('Hilton Worldwide')
            color_list[Hilton_index]='0000FF'
        except:
            None
        try:
            Hyatt_index=book_brand_list.tolist().index('Hyatt Hotels Corporation')
            color_list[Hyatt_index]='8968CD'
        except:
            None
        try:    
            Intercontinental_index=book_brand_list.tolist().index('Intercontinental Hotel Group')
            color_list[Intercontinental_index]='FF8000'
        except:
            None
        try:
            Marriott_index=book_brand_list.tolist().index('Marriott International')
            color_list[Marriott_index]='FF4500'
        except:
            None
        try:
            Wyndham_index=book_brand_list.tolist().index('Wyndham Hotel Group')
            color_list[Marriott_index]='8EE5EE'

        except:
            None
        try:
            Independent_index=book_brand_list.tolist().index('Independent')
            color_list[Marriott_index]='8B8B83'
        except:
            None
        


        brand_color = {}
        brand_color = {book_brand_list[i]: color_list[i] for i in range(len(book_brand_list))}

        for idx, point in enumerate(chart.series[0].points):

            col_idx = idx % len(color_list)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

    if brand_value_hotel.shape[0]>0:

        brand_value_hotel['brand_text']=brand_value_hotel['brand'] + " - $" +brand_value_hotel["total"].map('{:,.0f}'.format)

        brand_value_hotel=brand_value_hotel.sort_values(by='brand_text',ascending=True)

        new_brand_list=set(book_brand_list).intersection(brand_value_hotel['brand'])
        new_brand_list=list(new_brand_list)
        new_brand_list=sorted(new_brand_list)

        chart_data = ChartData()
        chart_data.categories =brand_value_hotel['brand_text']

        filtered_dictionary = {key: value for key, value in brand_color.items() if key in new_brand_list}
        new_colors=list(filtered_dictionary.values())

        chart_data.add_series('Contract Value Hotels', brand_value_hotel['total']/brand_value_hotel['total'].sum())

        x, y, cx, cy = Cm(18.3), Cm(2), Cm(15), Cm(16)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart


        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        for idx, point in enumerate(chart.series[0].points):

            col_idx = idx % len(new_colors)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor.from_string(new_colors[col_idx])

    left = Cm(5.11) 
    height = Cm(1.28)
    width = Cm(5.7)
    top=Cm(17.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		                          
    p.text = "Total Contracts {}".format(str(book_brand_hotel['id'].sum()))

    left = Cm(22.55) 
    height = Cm(1.28)
    width = Cm(5.7)
    top=Cm(17.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()		     
                        
    bhv=brand_value_hotel['total'].sum()
    p.text = "Total Spend ${}".format(str(f'{bhv:,.0f}' ))
    

    slide = prs.slides[7]

    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame


    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Other Summary for {}".format(q_text)

    font = run.font
    font.name = 'Courier New'
    font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(71,54,241)

    if book_brand_non_hotel.shape[0]>0:
        book_brand_non_hotel['brand_text']=book_brand_non_hotel['p_type'] + " - " + book_brand_non_hotel['id'].astype(str)

        book_brand_non_hotel=book_brand_non_hotel.sort_values(by='brand_text',ascending=True)

        chart_data = ChartData()
        chart_data.categories =book_brand_non_hotel['brand_text']

        chart_data.add_series('Booked Brand Distribution Other', book_brand_non_hotel['id']/book_brand_non_hotel['id'].sum())

        x, y, cx, cy = Cm(0.5), Cm(2), Cm(15), Cm(16)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'

        number_of_colors = book_brand_non_hotel.shape[0]
        import random

        color_list=[]
        color_list = [''.join([random.choice('0123456789ABCDEF') for j in range(6)]) for i in range(number_of_colors)]
        book_brand_list=book_brand_non_hotel['p_type'].unique()

        try:
            Event_index=book_brand_list.tolist().index('Event Venue')
            color_list[Event_index]='F08080'
        except:
            None
        try:
            Restaurant_index=book_brand_list.tolist().index('Restaurant')
            color_list[Restaurant_index]='BCD2EE'
        except:
            None
        try:    
            Unique_index=book_brand_list.tolist().index('Unique Venue')
            color_list[Unique_index]='BDFCC9'
        except:
            None
       
        brand_color = {}
        brand_color = {book_brand_list[i]: color_list[i] for i in range(len(book_brand_list))}


        for idx, point in enumerate(chart.series[0].points):

            col_idx = idx % len(color_list)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

    if brand_value_non_hotel.shape[0]>0:

        brand_value_non_hotel['brand_text']=brand_value_non_hotel['p_type'] + " - $" +brand_value_non_hotel["total"].map('{:,.0f}'.format)

        brand_value_non_hotel=brand_value_non_hotel.sort_values(by='brand_text',ascending=True)

        new_brand_list=set(book_brand_list).intersection(brand_value_non_hotel['p_type'])
        new_brand_list=list(new_brand_list)
        new_brand_list=sorted(new_brand_list)

        chart_data = ChartData()
        chart_data.categories =brand_value_non_hotel['brand_text']

        filtered_dictionary = {key: value for key, value in brand_color.items() if key in new_brand_list}
        new_colors=list(filtered_dictionary.values())

        chart_data.add_series('Contract Value Other', brand_value_non_hotel['total']/brand_value_non_hotel['total'].sum())

        x, y, cx, cy = Cm(18.3), Cm(2), Cm(15), Cm(16)
        chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        for idx, point in enumerate(chart.series[0].points):

            col_idx = idx % len(new_colors)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor.from_string(new_colors[col_idx])

    left = Cm(5.11) 
    height = Cm(1.28)
    width = Cm(5.7)
    top=Cm(17.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()                                
    p.text = "Total Contracts {}".format(str(book_brand_non_hotel['id'].sum()))

    left = Cm(22.55) 
    height = Cm(1.28)
    width = Cm(5.7)
    top=Cm(17.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)    # Adding Shape object (Text Box)

    tf = txBox.text_frame

    p = tf.add_paragraph()   
    bvnh=brand_value_non_hotel['total'].sum()                             
    p.text = "Total Spend ${}".format(str(f'{bvnh:,.0f}'))


    

    slide = prs.slides[8]
    
    left = Cm(1.04) 
    width = Cm(17.87)
    height= Cm(1.03)
    top=Cm(0.97)

    txBox = slide.shapes.add_textbox(left, top, width, height) 
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{} 2022 Reporting".format(q_text)

    font = run.font
    font.name = 'Courier New'
    #font.bold=True
    font.size = Pt(18)
    font.color.rgb = RGBColor(255,255,255)

    to_replace='{} Totals'.format(q_text)
    def replace_text(replacements, shapes):
        for shape in shapes:
            for match, replacement in replacements.items():
                if shape.has_text_frame:
                    if (shape.text.find(match)) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            whole_text = "".join(run.text for run in paragraph.runs)
                            whole_text = whole_text.replace(str(match), str(replacement))
                            for idx, run in enumerate(paragraph.runs):
                                if idx != 0:
                                    p = paragraph._p
                                    p.remove(run._r)
                            if bool(paragraph.runs):
                                paragraph.runs[0].text = whole_text

    slides = [slide for slide in prs.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)

    replaces = {
                        'Q2 Totals': to_replace
                }

    replace_text(replaces, shapes)

    prs.save(ppt_name)

""" ****************************************************************** CREATION OF PPTX SLIDE 1  *************************************************************************"""




    

